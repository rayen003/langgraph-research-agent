"""Phase B1 — deck theme resolution + themed assembly smoke tests."""

from __future__ import annotations

import utils as utils_mod
from graphs.workflows.deck.assemble import assemble_pptx_node
from graphs.workflows.deck.theme import DeckTheme, resolve_theme
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt


# ── resolve_theme ──────────────────────────────────────────────────────────


def test_audience_presets_distinct():
    board = resolve_theme({"audience": "board"})
    internal = resolve_theme({"audience": "internal"})
    ic = resolve_theme({"audience": "ic"})

    assert board.density == "spacious"
    assert internal.density == "compact"
    assert ic.density == "standard"
    # board is roomier + larger; internal denser + smaller.
    assert board.font_scale > ic.font_scale > internal.font_scale
    assert board.margin_x > internal.margin_x


def test_default_accent_is_teal():
    assert resolve_theme({}).accent == RGBColor(0x2D, 0xD4, 0xBF)


def test_brief_tokens_override_audience():
    t = resolve_theme(
        {"audience": "board", "density": "compact", "font_scale": 0.9, "accent": "#ff0000"}
    )
    assert t.density == "compact"          # overrides board's spacious
    assert t.font_scale == 0.9             # overrides board's 1.1
    assert t.accent == RGBColor(0xFF, 0x00, 0x00)


def test_font_scale_clamped():
    assert resolve_theme({"font_scale": 5.0}).font_scale == 1.4
    assert resolve_theme({"font_scale": 0.1}).font_scale == 0.8


def test_invalid_accent_falls_back():
    # Garbage accent → audience/default fallback, not a crash.
    assert resolve_theme({"accent": "not-a-color"}).accent == RGBColor(0x2D, 0xD4, 0xBF)
    assert resolve_theme({"audience": "board", "accent": "zzz"}).accent == RGBColor(0x1F, 0x3A, 0x5F)


def test_three_digit_hex_expands():
    assert resolve_theme({"accent": "#abc"}).accent == RGBColor(0xAA, 0xBB, 0xCC)


def test_unknown_density_defaults_standard():
    assert resolve_theme({"density": "ginormous"}).density == "standard"


def test_size_scales_and_space_tracks_density():
    t = resolve_theme({"font_scale": 1.0})
    assert t.size(20) == Pt(20)
    t2 = resolve_theme({"font_scale": 1.5})  # clamped to 1.4
    assert t2.size(20) == Pt(round(20 * 1.4))
    # Spacious puts more space after paragraphs than compact.
    assert resolve_theme({"density": "spacious"}).space(10) > resolve_theme(
        {"density": "compact"}
    ).space(10)


def test_resolve_theme_is_deterministic():
    brief = {"audience": "client", "font_scale": 1.05}
    assert resolve_theme(brief) == resolve_theme(dict(brief))
    assert isinstance(resolve_theme(brief), DeckTheme)


# ── themed assembly smoke ────────────────────────────────────────────────────


def _slides() -> list[dict]:
    return [
        {"slide_id": "s1", "layout": "title", "title": "AAPL Case",
         "body_paragraphs": ["Market expectations vs reality"]},
        {"slide_id": "s2", "layout": "bullets", "title": "Drivers",
         "body_bullets": ["Revenue growth", "Margin expansion", "Buybacks"]},
        {"slide_id": "s3", "layout": "three_box", "title": "Summary",
         "columns": [
             {"heading": "Priced", "bullets": ["8% growth"]},
             {"heading": "Assumed", "bullets": ["6% growth"]},
             {"heading": "Required", "bullets": ["10% growth"]},
         ]},
        {"slide_id": "s4", "layout": "references", "title": "Sources",
         "body_bullets": ["ev_1", "ev_2"]},
    ]


def test_assemble_renders_for_each_audience(tmp_path, monkeypatch):
    monkeypatch.setattr(utils_mod, "RUNS_DIR", tmp_path)
    utils_mod.set_thread_id("tdeck")

    for audience in ("board", "ic", "internal", "client", "generic"):
        state = {
            "slides": _slides(),
            "brief": {"title": f"Deck {audience}", "audience": audience},
            "parent_step_id": "wf",
        }
        out = assemble_pptx_node(state)
        path = out["pptx_path"]
        assert path.endswith(".pptx")
        prs = Presentation(path)
        assert len(prs.slides) == 4


def test_assemble_honors_explicit_tokens(tmp_path, monkeypatch):
    monkeypatch.setattr(utils_mod, "RUNS_DIR", tmp_path)
    utils_mod.set_thread_id("tdeck2")

    state = {
        "slides": _slides(),
        "brief": {
            "title": "Custom", "audience": "ic",
            "density": "spacious", "accent": "#ff8800", "font_scale": 1.2,
        },
        "parent_step_id": "wf",
    }
    out = assemble_pptx_node(state)
    assert Presentation(out["pptx_path"]).slides is not None

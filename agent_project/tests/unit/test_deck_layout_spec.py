"""Phase B2 — LLM-emitted declarative layout_spec validation + render smoke."""

from __future__ import annotations

import utils as utils_mod
from graphs.workflows.deck.assemble import assemble_pptx_node
from graphs.workflows.deck.state import LayoutRegion, LayoutSpec, SlideContent
from pptx import Presentation


# ── LayoutRegion.in_bounds ───────────────────────────────────────────────────


def test_region_in_bounds_accepts_valid():
    assert LayoutRegion(kind="text", x=0.1, y=0.1, w=0.5, h=0.3).in_bounds()


def test_region_rejects_zero_size():
    # Missing w/h default to 0 → degenerate → rejected (graceful, no crash).
    assert not LayoutRegion(kind="text", x=0.1, y=0.1).in_bounds()


def test_region_rejects_overflow():
    assert not LayoutRegion(kind="text", x=0.6, y=0.1, w=0.5, h=0.3).in_bounds()  # x+w > 1
    assert not LayoutRegion(kind="text", x=0.1, y=0.8, w=0.3, h=0.3).in_bounds()  # y+h > 1


def test_region_rejects_negative():
    assert not LayoutRegion(kind="text", x=-0.1, y=0.1, w=0.5, h=0.3).in_bounds()


# ── LayoutRegion.has_content ─────────────────────────────────────────────────


def test_has_content_per_kind():
    assert LayoutRegion(kind="text", text="hi").has_content()
    assert not LayoutRegion(kind="text", text="   ").has_content()
    assert LayoutRegion(kind="bullets", items=["a", ""]).has_content()
    assert not LayoutRegion(kind="bullets", items=["", "  "]).has_content()
    assert LayoutRegion(kind="table", rows=[["h"]]).has_content()
    assert not LayoutRegion(kind="table", rows=[]).has_content()
    assert LayoutRegion(kind="image", image_path="/x.png").has_content()
    assert not LayoutRegion(kind="image").has_content()
    assert LayoutRegion(kind="accent_bar").has_content()  # decoration always "content"


# ── LayoutSpec.is_renderable ─────────────────────────────────────────────────


def _ok_region() -> LayoutRegion:
    return LayoutRegion(kind="text", x=0.1, y=0.1, w=0.5, h=0.3, text="hello", role="body")


def test_spec_renderable_happy():
    assert LayoutSpec(regions=[_ok_region()]).is_renderable()


def test_spec_empty_not_renderable():
    assert not LayoutSpec(regions=[]).is_renderable()


def test_spec_too_many_regions_not_renderable():
    assert not LayoutSpec(regions=[_ok_region() for _ in range(13)]).is_renderable()


def test_spec_out_of_bounds_region_kills_spec():
    bad = LayoutRegion(kind="text", x=0.9, y=0.1, w=0.5, h=0.3, text="x")
    assert not LayoutSpec(regions=[_ok_region(), bad]).is_renderable()


def test_spec_accent_only_not_renderable():
    # An accent bar alone carries no real content → not renderable.
    bar = LayoutRegion(kind="accent_bar", x=0.1, y=0.1, w=0.4, h=0.05, fill="accent")
    assert not LayoutSpec(regions=[bar]).is_renderable()


# ── Themed assembly with a layout_spec slide ─────────────────────────────────


def test_assemble_renders_layout_spec_slide(tmp_path, monkeypatch):
    monkeypatch.setattr(utils_mod, "RUNS_DIR", tmp_path)
    utils_mod.set_thread_id("tspec")

    spec = LayoutSpec(regions=[
        LayoutRegion(kind="accent_bar", x=0.06, y=0.06, w=0.3, h=0.02, fill="accent"),
        LayoutRegion(kind="text", x=0.06, y=0.10, w=0.88, h=0.12, text="Custom Title", role="title"),
        LayoutRegion(kind="text", text="$240", x=0.06, y=0.3, w=0.4, h=0.2, role="metric"),
        LayoutRegion(kind="bullets", x=0.5, y=0.3, w=0.44, h=0.4,
                     items=["Point one", "Point two"], role="body"),
        LayoutRegion(kind="table", x=0.06, y=0.6, w=0.88, h=0.25,
                     rows=[["Metric", "Value"], ["WACC", "8.1%"]]),
    ])
    assert spec.is_renderable()  # must exercise the _render_spec path, not fallback
    sc = SlideContent(
        slide_id="s1", layout="bullets", title="Custom Title",
        body_bullets=["fallback bullet"],  # canonical content still present
        layout_spec=spec,
    )
    state = {
        "slides": [sc.model_dump()],
        "brief": {"title": "Spec Deck", "audience": "ic"},
        "parent_step_id": "wf",
    }
    out = assemble_pptx_node(state)
    prs = Presentation(out["pptx_path"])
    assert len(prs.slides) == 1


def test_assemble_falls_back_when_spec_invalid(tmp_path, monkeypatch):
    monkeypatch.setattr(utils_mod, "RUNS_DIR", tmp_path)
    utils_mod.set_thread_id("tspec2")

    # Overflowing region → spec not renderable → standard bullets renderer used.
    bad_spec = LayoutSpec(regions=[
        LayoutRegion(kind="text", x=0.9, y=0.1, w=0.5, h=0.3, text="overflow"),
    ])
    assert not bad_spec.is_renderable()
    sc = SlideContent(
        slide_id="s1", layout="bullets", title="Fallback",
        body_bullets=["a", "b", "c"], layout_spec=bad_spec,
    )
    state = {
        "slides": [sc.model_dump()],
        "brief": {"title": "Fallback Deck", "audience": "ic"},
        "parent_step_id": "wf",
    }
    out = assemble_pptx_node(state)
    assert len(Presentation(out["pptx_path"]).slides) == 1

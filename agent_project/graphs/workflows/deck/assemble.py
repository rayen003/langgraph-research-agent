"""PPTX assembly — render ``SlideContent`` list to ``.pptx``.

Uses ``python-pptx`` directly (no template, all programmatic).  One small
renderer function per layout keeps each render simple and testable.  Future
template-library work plugs in here: layout dispatch checks for a cached
template before falling through to the programmatic renderer.

Visual style is driven by a ``DeckTheme`` (see ``theme.py``), resolved once
per run from the brief's audience + optional ``density`` / ``accent`` /
``font_scale`` tokens, then threaded into every renderer as ``t``.  Colors,
font sizes (via ``t.size``), paragraph spacing (via ``t.space``) and page
margins all come from the theme — no module-level design constants.

Slide design (deliberately minimal — looks clean, doesn't ape McKinsey):
  - 16:9 widescreen (13.33 x 7.5 inches).
  - No background images, no shadows, no clutter.

Output: file at ``runs/<thread_id>/decks/<sanitized_title>.pptx``.
"""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from utils import get_run_dir

from .activity import emit_step
from .state import DeckState, LayoutRegion, LayoutSpec, SlideContent
from .theme import DeckTheme, resolve_theme

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Fixed page geometry (16:9). Margins + all style tokens live on DeckTheme.
# ---------------------------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------


def assemble_pptx_node(state: DeckState) -> dict:
    """Render the slides list to a PPTX file on disk."""
    parent_step_id = state.get("parent_step_id") or "workflow_deck"
    emit_step("assemble_pptx", "start", parent_step_id)

    slides_raw = state.get("slides") or []
    if not slides_raw:
        emit_step("assemble_pptx", "error", parent_step_id, {
            "summary_line": "No slides to assemble.",
        })
        raise ValueError("assemble_pptx_node received empty slides list.")

    brief = state.get("brief") or {}
    deck_title = str(brief.get("title") or "deck")
    theme = resolve_theme(brief)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    layout_blank = prs.slide_layouts[6]  # blank layout — fully manual placement

    for i, raw in enumerate(slides_raw):
        try:
            sc = SlideContent.model_validate(raw)
        except Exception as exc:  # noqa: BLE001
            logger.warning("Skipping malformed slide #%d: %s", i, exc)
            continue
        slide = prs.slides.add_slide(layout_blank)
        _render(slide, sc, theme, deck_index=i + 1, deck_total=len(slides_raw))

    # Write to disk under run dir.
    run_dir = Path(get_run_dir())
    decks_dir = run_dir / "decks"
    decks_dir.mkdir(parents=True, exist_ok=True)
    safe_name = _sanitize_filename(deck_title) or "deck"
    pptx_path = decks_dir / f"{safe_name}.pptx"
    prs.save(str(pptx_path))

    emit_step("assemble_pptx", "complete", parent_step_id, {
        "summary_line": f"Wrote {pptx_path.name} ({len(prs.slides)} slides)",
        "pptx_path": str(pptx_path),
        "slide_count": len(prs.slides),
    })
    logger.info(
        "Deck assembled at %s (%d slides, audience=%s density=%s scale=%.2f)",
        pptx_path, len(prs.slides), brief.get("audience"), theme.density, theme.font_scale,
    )
    return {"pptx_path": str(pptx_path)}


# ---------------------------------------------------------------------------
# Layout dispatch
# ---------------------------------------------------------------------------


def _render(
    slide, content: SlideContent, t: DeckTheme, *, deck_index: int, deck_total: int
) -> None:
    # Phase B2: a valid LLM-emitted layout_spec wins over the canned layout.
    spec = content.layout_spec
    if spec is not None and spec.is_renderable():
        try:
            _render_spec(slide, spec, t)
            _render_footer(slide, content, t, deck_index=deck_index, deck_total=deck_total)
            return
        except Exception:
            logger.exception(
                "layout_spec render failed for slide_id=%s — falling back to layout=%s.",
                content.slide_id, content.layout,
            )

    renderer = _LAYOUT_RENDERERS.get(content.layout, _render_bullets)
    try:
        renderer(slide, content, t)
    except Exception:
        logger.exception("Renderer failed for layout=%s slide_id=%s — falling back to bullets.",
                         content.layout, content.slide_id)
        _render_bullets(slide, content, t)
    _render_footer(slide, content, t, deck_index=deck_index, deck_total=deck_total)


# ---------------------------------------------------------------------------
# Phase B2 — generic layout_spec interpreter
# ---------------------------------------------------------------------------

# Role → (base font pt, color attr name, bold). Sizes scale via DeckTheme.size.
_ROLE_STYLE: dict[str, tuple[float, str, bool]] = {
    "title":   (44, "dark", True),
    "heading": (20, "dark", True),
    "body":    (16, "body", False),
    "caption": (11, "muted", False),
    "metric":  (60, "accent", True),
    "muted":   (13, "muted", False),
}

_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}


def _role_style(t: DeckTheme, role: str | None) -> tuple[Pt, RGBColor, bool]:
    base_pt, color_attr, bold = _ROLE_STYLE.get(role or "body", _ROLE_STYLE["body"])
    return t.size(base_pt), getattr(t, color_attr), bold


def _render_spec(slide, spec: LayoutSpec, t: DeckTheme) -> None:
    """Paint each region of a validated ``LayoutSpec`` at fractional coords."""
    for r in spec.regions:
        left = Emu(int(int(SLIDE_W) * r.x))
        top = Emu(int(int(SLIDE_H) * r.y))
        width = Emu(int(int(SLIDE_W) * r.w))
        height = Emu(int(int(SLIDE_H) * r.h))
        if r.kind == "accent_bar":
            _spec_accent_bar(slide, t, r, left, top, width, height)
        elif r.kind == "image":
            _spec_image(slide, t, r, left, top, width, height)
        elif r.kind == "table":
            _spec_table(slide, t, r, left, top, width, height)
        else:  # text | bullets
            _spec_text(slide, t, r, left, top, width, height)


def _spec_text(slide, t: DeckTheme, r: LayoutRegion, left, top, width, height) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    size, color, bold = _role_style(t, r.role)
    align = _ALIGN.get(r.align or "")
    if r.kind == "bullets":
        items = [str(i) for i in r.items if str(i).strip()]
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"•  {item}"
            if align is not None:
                p.alignment = align
            for run in p.runs:
                run.font.size = size
                run.font.color.rgb = color
                run.font.bold = bold
            p.space_after = t.space(6)
    else:
        _set_text(tf, r.text or "", size=size, bold=bold, color=color, align=align)


def _spec_table(slide, t: DeckTheme, r: LayoutRegion, left, top, width, height) -> None:
    rows = [row for row in r.rows if row]
    if not rows:
        return
    n_rows = len(rows)
    n_cols = max(len(row) for row in rows)
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    for ri, row in enumerate(rows):
        for ci in range(n_cols):
            cell = table.cell(ri, ci)
            cell.text = str(row[ci]) if ci < len(row) else ""
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = t.size(12)
                    run.font.color.rgb = t.dark if ri == 0 else t.body
                    if ri == 0:
                        run.font.bold = True


def _spec_image(slide, t: DeckTheme, r: LayoutRegion, left, top, width, height) -> None:
    if r.image_path and Path(r.image_path).exists():
        slide.shapes.add_picture(r.image_path, left, top, width=width, height=height)
    else:
        box = slide.shapes.add_textbox(left, top, width, height)
        _set_text(box.text_frame, "(image unavailable)", size=t.size(14),
                  color=t.muted, align=PP_ALIGN.CENTER)


def _spec_accent_bar(slide, t: DeckTheme, r: LayoutRegion, left, top, width, height) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = t.panel_fill if r.fill == "panel" else t.accent
    bar.line.fill.background()
    bar.shadow.inherit = False


# ---------------------------------------------------------------------------
# Per-layout renderers (intentionally small — easy to swap for templates)
# ---------------------------------------------------------------------------


def _render_title(slide, c: SlideContent, t: DeckTheme) -> None:
    """Centered title + optional subtitle from first paragraph."""
    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, t.margin_x, Inches(3.0), Inches(0.6), Inches(0.08),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = t.accent
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(t.margin_x, Inches(3.2), SLIDE_W - 2 * t.margin_x, Inches(1.5))
    _set_text(title_box.text_frame, c.title, size=t.size(44), bold=True, color=t.dark)

    if c.body_paragraphs:
        subtitle_box = slide.shapes.add_textbox(t.margin_x, Inches(4.6), SLIDE_W - 2 * t.margin_x, Inches(1.2))
        _set_text(subtitle_box.text_frame, c.body_paragraphs[0], size=t.size(18), color=t.muted)


def _render_section_header(slide, c: SlideContent, t: DeckTheme) -> None:
    """Big centered section divider."""
    title_box = slide.shapes.add_textbox(t.margin_x, Inches(3.0), SLIDE_W - 2 * t.margin_x, Inches(1.5))
    tf = title_box.text_frame
    _set_text(tf, c.title, size=t.size(40), bold=True, color=t.dark, align=PP_ALIGN.CENTER)


def _render_bullets(slide, c: SlideContent, t: DeckTheme) -> None:
    """Title + bulleted list."""
    _render_title_strip(slide, c.title, t)
    body_top = Inches(1.6)
    body_h = SLIDE_H - body_top - Inches(0.6)
    body_box = slide.shapes.add_textbox(t.margin_x, body_top, SLIDE_W - 2 * t.margin_x, body_h)
    tf = body_box.text_frame
    tf.word_wrap = True
    bullets = c.body_bullets or [p[:140] for p in c.body_paragraphs[:6]]
    if not bullets:
        bullets = ["(No content)"]
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {b}"
        for run in p.runs:
            run.font.size = t.size(16)
            run.font.color.rgb = t.body
        p.space_after = t.space(8)


def _render_narrative(slide, c: SlideContent, t: DeckTheme) -> None:
    """Title + 1-3 paragraphs."""
    _render_title_strip(slide, c.title, t)
    body_box = slide.shapes.add_textbox(t.margin_x, Inches(1.6), SLIDE_W - 2 * t.margin_x, Inches(5.2))
    tf = body_box.text_frame
    tf.word_wrap = True
    # Operator precedence fix: original `A or B if cond else C` parses as
    # `(A or B) if cond else C`, so when body_paragraphs HAS content but
    # body_bullets is empty the whole expression falls into the else and
    # produced "(No content)" — that wiped every thesis/narrative slide.
    if c.body_paragraphs:
        paras = c.body_paragraphs
    elif c.body_bullets:
        paras = [" ".join(c.body_bullets)]
    else:
        paras = ["(No content)"]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = para
        for run in p.runs:
            run.font.size = t.size(18)
            run.font.color.rgb = t.body
        p.space_after = t.space(12)


def _render_metric_callout(slide, c: SlideContent, t: DeckTheme) -> None:
    """Title + huge accent number derived from first paragraph + supporting text."""
    _render_title_strip(slide, c.title, t)
    # Big number: try to extract a $X.XX pattern; else use first 12 chars of paragraph.
    para = (c.body_paragraphs or [""])[0]
    m = re.search(r"\$[\d,]+(?:\.\d+)?(?:[MBKT]|%|bps)?", para)
    headline = m.group(0) if m else para[:20]
    big = slide.shapes.add_textbox(t.margin_x, Inches(2.0), SLIDE_W - 2 * t.margin_x, Inches(2.0))
    _set_text(big.text_frame, headline, size=t.size(72), bold=True, color=t.accent)

    if para:
        sub = slide.shapes.add_textbox(t.margin_x, Inches(4.5), SLIDE_W - 2 * t.margin_x, Inches(2.0))
        _set_text(sub.text_frame, para, size=t.size(18), color=t.body)


def _render_scenario_table(slide, c: SlideContent, t: DeckTheme) -> None:
    """Title + table from c.table_rows."""
    _render_title_strip(slide, c.title, t)
    rows = c.table_rows
    if not rows:
        _render_bullets(slide, c, t)
        return
    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    table_top = Inches(1.6)
    table_h = min(SLIDE_H - table_top - Inches(0.8), Inches(0.5) * n_rows)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, t.margin_x, table_top, SLIDE_W - 2 * t.margin_x, table_h,
    )
    table = table_shape.table
    for ri, row in enumerate(rows):
        for ci in range(n_cols):
            cell = table.cell(ri, ci)
            text = str(row[ci]) if ci < len(row) else ""
            cell.text = text
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = t.size(12)
                    run.font.color.rgb = t.dark if ri == 0 else t.body
                    if ri == 0:
                        run.font.bold = True


def _render_chart_caption(slide, c: SlideContent, t: DeckTheme) -> None:
    """Title + chart image + caption."""
    _render_title_strip(slide, c.title, t)
    chart_h = Inches(4.5)
    chart_w = Inches(8.0)
    chart_top = Inches(1.8)
    chart_left = Emu(int((SLIDE_W - chart_w) / 2))

    if c.chart_path and Path(c.chart_path).exists():
        slide.shapes.add_picture(c.chart_path, chart_left, chart_top, width=chart_w, height=chart_h)
    else:
        placeholder = slide.shapes.add_textbox(chart_left, chart_top, chart_w, chart_h)
        _set_text(placeholder.text_frame, "(chart unavailable)", size=t.size(18), color=t.muted, align=PP_ALIGN.CENTER)

    if c.chart_caption or c.body_paragraphs:
        cap = c.chart_caption or c.body_paragraphs[0]
        cap_box = slide.shapes.add_textbox(t.margin_x, Inches(6.5), SLIDE_W - 2 * t.margin_x, Inches(0.6))
        _set_text(cap_box.text_frame, cap, size=t.size(12), color=t.muted, align=PP_ALIGN.CENTER)


def _render_references(slide, c: SlideContent, t: DeckTheme) -> None:
    """References slide — title + bulleted citation list."""
    _render_title_strip(slide, c.title or "Sources", t)
    body_box = slide.shapes.add_textbox(t.margin_x, Inches(1.6), SLIDE_W - 2 * t.margin_x, Inches(5.2))
    tf = body_box.text_frame
    tf.word_wrap = True
    bullets = c.body_bullets or ["(no sources)"]
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        for run in p.runs:
            run.font.size = t.size(11)
            run.font.color.rgb = t.muted


# ---------------------------------------------------------------------------
# Expectations-deck layouts (institutional shape)
# ---------------------------------------------------------------------------


def _render_reconciliation_table(slide, c: SlideContent, t: DeckTheme) -> None:
    """Reconciliation table: title + tight comparison table + leader paragraph.

    Expected ``c.table_rows`` shape (LLM-produced, first row = header):
        [["Metric", "Model", "Market-Implied", "Gap"], …]
    Falls back to ``_render_scenario_table`` styling when populated; if the
    LLM only filled ``body_paragraphs`` the leader text still appears.
    """
    _render_title_strip(slide, c.title, t)

    # Optional leader paragraph above the table (the "headline takeaway").
    table_top = Inches(1.6)
    leader_paragraphs = c.body_paragraphs[:1]
    if leader_paragraphs:
        leader_box = slide.shapes.add_textbox(t.margin_x, table_top, SLIDE_W - 2 * t.margin_x, Inches(0.7))
        _set_text(leader_box.text_frame, leader_paragraphs[0], size=t.size(18), color=t.body)
        table_top = Inches(2.4)

    rows = c.table_rows
    if not rows:
        # No table → degrade to a bullet list of the remaining paragraphs.
        body_box = slide.shapes.add_textbox(t.margin_x, table_top, SLIDE_W - 2 * t.margin_x, Inches(4.5))
        tf = body_box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(c.body_paragraphs[1:5] or c.body_bullets[:5] or ["(no comparison data)"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"•  {line}"
            for run in p.runs:
                run.font.size = t.size(18)
                run.font.color.rgb = t.body
        return

    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    table_h = min(SLIDE_H - table_top - Inches(0.8), Inches(0.55) * n_rows)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, t.margin_x, table_top, SLIDE_W - 2 * t.margin_x, table_h,
    )
    table = table_shape.table
    for ri, row in enumerate(rows):
        for ci in range(n_cols):
            cell = table.cell(ri, ci)
            text = str(row[ci]) if ci < len(row) else ""
            cell.text = text
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = t.size(13)
                    run.font.color.rgb = t.dark if ri == 0 else t.body
                    if ri == 0:
                        run.font.bold = True


def _render_three_box(slide, c: SlideContent, t: DeckTheme) -> None:
    """Three vertical columns: Priced / Assumed / Required.

    Reads ``c.columns`` — list of {"heading": str, "bullets": [str], "paragraphs": [str]}.
    Falls back to splitting ``body_bullets`` into thirds when columns missing.
    """
    _render_title_strip(slide, c.title, t)

    cols = c.columns or []
    if not cols:
        # Synthesize from bullets if LLM failed to emit columns.
        n = max(1, len(c.body_bullets) // 3)
        chunks = [c.body_bullets[i:i + n] for i in range(0, len(c.body_bullets), n)][:3]
        cols = [
            {"heading": h, "bullets": chunks[i] if i < len(chunks) else []}
            for i, h in enumerate(["Priced", "Assumed", "Required"])
        ]
    cols = cols[:3]  # cap at three

    n = len(cols) or 1
    gutter = Inches(0.2)
    total_w = SLIDE_W - 2 * t.margin_x - gutter * (n - 1)
    col_w = Emu(int(total_w / n))
    col_top = Inches(1.7)
    col_h = SLIDE_H - col_top - Inches(0.6)

    for i, col in enumerate(cols):
        left = Emu(int(t.margin_x + (col_w + gutter) * i))

        # Background panel — light tint for hierarchy.
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, col_top, col_w, col_h)
        panel.fill.solid()
        panel.fill.fore_color.rgb = t.panel_fill
        panel.line.color.rgb = t.panel_line
        panel.shadow.inherit = False

        inner_x = Emu(int(left + Inches(0.25)))
        inner_w = Emu(int(col_w - Inches(0.5)))

        # Heading
        heading_box = slide.shapes.add_textbox(inner_x, Emu(int(col_top + Inches(0.2))), inner_w, Inches(0.6))
        _set_text(
            heading_box.text_frame,
            str(col.get("heading") or ""),
            size=t.size(18), bold=True, color=t.dark,
        )

        # Bullets
        bullets = col.get("bullets") or []
        body_box = slide.shapes.add_textbox(
            inner_x, Emu(int(col_top + Inches(0.9))), inner_w, Emu(int(col_h - Inches(1.1))),
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        for j, b in enumerate(bullets[:8]):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"•  {b}"
            for run in p.runs:
                run.font.size = t.size(13)
                run.font.color.rgb = t.body
            p.space_after = t.space(6)


def _render_two_col_narrative(slide, c: SlideContent, t: DeckTheme) -> None:
    """Two-column narrative: Bull vs Bear (or any binary frame).

    Reads ``c.columns`` (first 2 entries). Each column gets a heading and
    1-3 paragraphs. Falls back to splitting ``body_paragraphs`` in half.
    """
    _render_title_strip(slide, c.title, t)

    cols = c.columns or []
    if not cols and c.body_paragraphs:
        mid = max(1, len(c.body_paragraphs) // 2)
        cols = [
            {"heading": "Bull View", "paragraphs": c.body_paragraphs[:mid]},
            {"heading": "Bear View", "paragraphs": c.body_paragraphs[mid:]},
        ]
    cols = cols[:2]

    n = len(cols) or 1
    gutter = Inches(0.3)
    total_w = SLIDE_W - 2 * t.margin_x - gutter * (n - 1)
    col_w = Emu(int(total_w / n))
    col_top = Inches(1.7)
    col_h = SLIDE_H - col_top - Inches(0.6)

    accent_colors = [t.accent, t.accent_alt]  # bull / bear

    for i, col in enumerate(cols):
        left = Emu(int(t.margin_x + (col_w + gutter) * i))

        # Top accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, col_top, col_w, Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_colors[i % len(accent_colors)]
        bar.line.fill.background()

        # Heading
        heading_box = slide.shapes.add_textbox(left, Emu(int(col_top + Inches(0.2))), col_w, Inches(0.5))
        _set_text(
            heading_box.text_frame,
            str(col.get("heading") or ""),
            size=t.size(18), bold=True, color=t.dark,
        )

        # Paragraphs (or bullets)
        body_box = slide.shapes.add_textbox(
            left, Emu(int(col_top + Inches(0.9))), col_w, Emu(int(col_h - Inches(1.0))),
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        items = col.get("paragraphs") or col.get("bullets") or []
        for j, item in enumerate(items[:5]):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = str(item)
            for run in p.runs:
                run.font.size = t.size(14)
                run.font.color.rgb = t.body
            p.space_after = t.space(10)


def _render_flow_diagram(slide, c: SlideContent, t: DeckTheme) -> None:
    """Vertical flow diagram: ordered steps with arrows between them.

    Reads ``c.flow_steps`` — list of {"label": str, "detail": str}.
    Caption from ``c.body_paragraphs[0]`` if present.
    """
    _render_title_strip(slide, c.title, t)

    steps = c.flow_steps or []
    if not steps and c.body_bullets:
        # Fallback: treat each bullet as a {"label": bullet, "detail": ""} step.
        steps = [{"label": b, "detail": ""} for b in c.body_bullets[:6]]
    steps = steps[:6]

    if not steps:
        body = slide.shapes.add_textbox(t.margin_x, Inches(2.0), SLIDE_W - 2 * t.margin_x, Inches(2.0))
        _set_text(body.text_frame, "(no flow data)", size=t.size(18), color=t.muted, align=PP_ALIGN.CENTER)
        return

    n = len(steps)
    region_top = Inches(1.7)
    region_h = SLIDE_H - region_top - Inches(1.2)
    step_h = Emu(int(region_h / (n + (n - 1) * 0.3)))  # arrows occupy 30% of step height between nodes
    arrow_h = Emu(int(step_h * 0.3))
    step_w = Inches(7.0)
    step_left = Emu(int((SLIDE_W - step_w) / 2))

    for i, step in enumerate(steps):
        top = Emu(int(region_top + (step_h + arrow_h) * i))

        # Node
        node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, step_left, top, step_w, step_h)
        node.fill.solid()
        node.fill.fore_color.rgb = t.panel_fill if i % 2 == 0 else RGBColor(0xEC, 0xF1, 0xF7)
        node.line.color.rgb = t.accent
        node.shadow.inherit = False

        # Label + detail
        text_box = slide.shapes.add_textbox(
            Emu(int(step_left + Inches(0.3))), Emu(int(top + Inches(0.1))),
            Emu(int(step_w - Inches(0.6))), Emu(int(step_h - Inches(0.2))),
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        label = str(step.get("label") or "")
        detail = str(step.get("detail") or "")
        tf.text = label
        for run in tf.paragraphs[0].runs:
            run.font.size = t.size(15)
            run.font.bold = True
            run.font.color.rgb = t.dark
        if detail:
            p = tf.add_paragraph()
            p.text = detail
            for run in p.runs:
                run.font.size = t.size(12)
                run.font.color.rgb = t.muted

        # Down arrow to next node
        if i < n - 1:
            arrow_top = Emu(int(top + step_h))
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Emu(int(step_left + step_w / 2 - Inches(0.3))), arrow_top,
                Inches(0.6), arrow_h,
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = t.accent
            arrow.line.fill.background()


def _render_variable_impact_table(slide, c: SlideContent, t: DeckTheme) -> None:
    """Variable impact table: Variable | Δ | Impact %.

    Same as scenario_table but with bolded impact column. Caveat printed
    in muted text below table from ``c.chart_caption`` or last paragraph.
    """
    _render_title_strip(slide, c.title, t)
    rows = c.table_rows
    if not rows:
        body = slide.shapes.add_textbox(t.margin_x, Inches(2.0), SLIDE_W - 2 * t.margin_x, Inches(2.0))
        _set_text(body.text_frame, "(no impact rows)", size=t.size(18), color=t.muted)
        return

    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    table_top = Inches(1.7)
    table_h = min(SLIDE_H - table_top - Inches(1.4), Inches(0.55) * n_rows)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, t.margin_x, table_top, SLIDE_W - 2 * t.margin_x, table_h,
    )
    table = table_shape.table
    for ri, row in enumerate(rows):
        for ci in range(n_cols):
            cell = table.cell(ri, ci)
            text = str(row[ci]) if ci < len(row) else ""
            cell.text = text
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = t.size(13)
                    if ri == 0:
                        run.font.bold = True
                        run.font.color.rgb = t.dark
                    elif ci == n_cols - 1:
                        # Impact column: bold + accent
                        run.font.bold = True
                        run.font.color.rgb = t.accent
                    else:
                        run.font.color.rgb = t.body

    # Caveat line
    caveat = c.chart_caption or (c.body_paragraphs[-1] if c.body_paragraphs else "")
    if caveat:
        cap_box = slide.shapes.add_textbox(
            t.margin_x, SLIDE_H - Inches(0.9), SLIDE_W - 2 * t.margin_x, Inches(0.5),
        )
        _set_text(cap_box.text_frame, caveat, size=t.size(11), color=t.muted)


def _render_decision_summary(slide, c: SlideContent, t: DeckTheme) -> None:
    """Closing decision slide: lead paragraph + required-conditions list.

    Lead paragraph from ``c.body_paragraphs[0]``. Conditions from
    ``c.body_bullets`` (each bullet = one required condition).
    """
    _render_title_strip(slide, c.title, t)

    # Lead paragraph (the framing statement).
    lead = (c.body_paragraphs[0] if c.body_paragraphs else "").strip()
    body_top = Inches(1.7)
    if lead:
        lead_box = slide.shapes.add_textbox(t.margin_x, body_top, SLIDE_W - 2 * t.margin_x, Inches(1.2))
        _set_text(lead_box.text_frame, lead, size=t.size(18), bold=True, color=t.dark)
        body_top = Inches(3.0)

    # "Required conditions" sub-heading
    sub_box = slide.shapes.add_textbox(t.margin_x, body_top, SLIDE_W - 2 * t.margin_x, Inches(0.4))
    _set_text(sub_box.text_frame, "Required conditions for upside:", size=t.size(14), bold=True, color=t.accent)

    bullets = c.body_bullets or c.body_paragraphs[1:6]
    if not bullets:
        return
    list_box = slide.shapes.add_textbox(
        t.margin_x, Emu(int(body_top + Inches(0.5))), SLIDE_W - 2 * t.margin_x, Inches(3.0),
    )
    tf = list_box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets[:5]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {b}"
        for run in p.runs:
            run.font.size = t.size(15)
            run.font.color.rgb = t.body
        p.space_after = t.space(8)


# Layout name → renderer fn
_LAYOUT_RENDERERS = {
    "title":             _render_title,
    "section_header":    _render_section_header,
    "bullets":           _render_bullets,
    "metric_callout":    _render_metric_callout,
    "narrative":         _render_narrative,
    "thesis":            _render_narrative,           # thesis ≈ narrative styling
    "scenario_table":    _render_scenario_table,
    "risk_summary":      _render_bullets,
    "chart_caption":     _render_chart_caption,
    "executive_summary": _render_narrative,
    "references":        _render_references,
    # ── Expectations-deck layouts ─────────────────────────────────────────
    "reconciliation_table":  _render_reconciliation_table,
    "three_box":             _render_three_box,
    "two_col_narrative":     _render_two_col_narrative,
    "flow_diagram":          _render_flow_diagram,
    "variable_impact_table": _render_variable_impact_table,
    "decision_summary":      _render_decision_summary,
}


# ---------------------------------------------------------------------------
# Shared parts (title strip + footer)
# ---------------------------------------------------------------------------


def _render_title_strip(slide, title: str, t: DeckTheme) -> None:
    """Top title strip: accent bar + title text."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, t.margin_x, Inches(0.6), Inches(0.4), Inches(0.06),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = t.accent
    bar.line.fill.background()

    box = slide.shapes.add_textbox(t.margin_x, Inches(0.75), SLIDE_W - 2 * t.margin_x, Inches(0.8))
    _set_text(box.text_frame, title, size=t.size(32), bold=True, color=t.dark)


def _render_footer(
    slide, c: SlideContent, t: DeckTheme, *, deck_index: int, deck_total: int
) -> None:
    """Footer: slide number only, muted.

    Raw citation IDs are deliberately NOT printed here — they leaked KG IDs
    (``ev_fmp_…``, block IDs) onto every slide. Readable, resolved citations
    live on the dedicated references ("Sources") slide instead.
    """
    foot = slide.shapes.add_textbox(t.margin_x, SLIDE_H - Inches(0.4), SLIDE_W - 2 * t.margin_x, Inches(0.3))
    _set_text(foot.text_frame, f"{deck_index} / {deck_total}", size=t.size(9), color=t.muted)


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------


def _set_text(
    tf,
    text: str,
    *,
    size: Pt | None = None,
    bold: bool = False,
    color: RGBColor | None = None,
    align: Any | None = None,
) -> None:
    """Apply text + uniform run styling to the first paragraph of a frame."""
    tf.text = text
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    for run in p.runs:
        if size is not None:
            run.font.size = size
        if bold:
            run.font.bold = True
        if color is not None:
            run.font.color.rgb = color


def _sanitize_filename(name: str) -> str:
    """Strip filesystem-hostile chars from a deck title for use as filename."""
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", name.strip())
    return cleaned[:80].strip("._")


__all__ = ["assemble_pptx_node"]
